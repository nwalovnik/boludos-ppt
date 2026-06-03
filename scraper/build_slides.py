"""
Generador de PPTX con las publicaciones de la semana, estilo editorial Comex.

Cada slide incluye:
  - Banner cyan top
  - Título grande en cyan (mayúsculas)
  - Lede editorial (párrafo principal)
  - Bullets con narrativa específica por serie
  - Gráfico matplotlib (línea o barras) con últimos 24 meses
  - Tabla con detalle de los últimos 6 meses
  - Footer "Fuente: elaboración propia en base a..."

Output: publicaciones_semana.pptx en la raíz del repo.
"""
from __future__ import annotations

import io
import json
from datetime import datetime, timezone, timedelta
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).resolve().parent.parent
DATA_PATH = ROOT / "data.json"
OUT_PATH = ROOT / "publicaciones_semana.pptx"

# Paleta editorial Comex
GRIS_TX = RGBColor(0x3F, 0x3F, 0x3F)
GRIS_LX = RGBColor(0x6A, 0x67, 0x60)
GRIS_BG = RGBColor(0xF2, 0xF2, 0xF0)
CYAN    = RGBColor(0x00, 0xB2, 0xC9)
CYAN_HX = "#00B2C9"
ROJO    = RGBColor(0xB9, 0x1C, 0x1C)
ROJO_HX = "#B91C1C"
VERDE   = RGBColor(0x15, 0x80, 0x3D)
VERDE_HX= "#15803D"
NEGRO   = RGBColor(0x1A, 0x1A, 0x1A)
BLANCO  = RGBColor(0xFF, 0xFF, 0xFF)
FONT    = "Encode Sans"

MESES_ES = ["enero","febrero","marzo","abril","mayo","junio",
            "julio","agosto","septiembre","octubre","noviembre","diciembre"]

# ──────────────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────────────
def lbl_mes(periodo: str) -> str:
    if not periodo or len(periodo) < 7 or "-" not in periodo:
        return str(periodo)
    y, m = periodo[:4], int(periodo[5:7])
    return f"{MESES_ES[m-1]} {y}"

def fmt_n(v, decimals=0) -> str:
    if v is None or v == "":
        return "—"
    if isinstance(v, (int, float)):
        s = f"{v:,.{decimals}f}"
        return s.replace(",", "X").replace(".", ",").replace("X", ".")
    return str(v)

def fmt_pct(v, decimals=1) -> str:
    if v is None: return "—"
    sign = "+" if v >= 0 else ""
    return f"{sign}{v:.{decimals}f}%"

def fmt_money(v, scale=1, unit="M$") -> str:
    if v is None: return "—"
    return f"{fmt_n(v/scale, 1)} {unit}"

def ultimo_lunes() -> datetime:
    d = datetime.now(timezone.utc).replace(hour=0, minute=0, second=0, microsecond=0)
    return d - timedelta(days=d.weekday())

# ──────────────────────────────────────────────────────────────────────────────
# Matplotlib chart builder con estilo editorial Comex
# ──────────────────────────────────────────────────────────────────────────────
def style_axes(ax, title=None):
    """Aplica el estilo editorial Comex a un Axes matplotlib."""
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#9A9990")
    ax.spines["bottom"].set_color("#9A9990")
    ax.spines["left"].set_linewidth(0.7)
    ax.spines["bottom"].set_linewidth(0.7)
    ax.tick_params(axis="x", rotation=45, labelsize=8.5, colors="#525050", length=2.5)
    ax.tick_params(axis="y", labelsize=8.5, colors="#525050", length=2.5)
    ax.grid(True, axis="y", alpha=0.18, linewidth=0.5)
    if title:
        ax.set_title(title, fontsize=11, color="#1A1A1A", loc="left",
                     pad=10, fontweight="bold")

def line_chart(arr, x_key, y_keys, labels, colors, title=None, y_fmt=None, n_meses=24):
    """Chart de líneas con últimos N meses. Devuelve BytesIO con PNG."""
    if not arr:
        return None
    data = arr[-n_meses:]
    x = [r.get(x_key, "") for r in data]
    fig, ax = plt.subplots(figsize=(6.5, 3.4), dpi=140)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    for y_key, lbl, col in zip(y_keys, labels, colors):
        y = [r.get(y_key) for r in data]
        ax.plot(x, y, color=col, linewidth=2.0, marker="o", markersize=3.2,
                label=lbl, markerfacecolor=col, markeredgecolor="white", markeredgewidth=0.6)
    style_axes(ax, title)
    if y_fmt:
        ax.yaxis.set_major_formatter(FuncFormatter(y_fmt))
    if len(labels) > 1:
        ax.legend(loc="best", fontsize=8.5, frameon=False, labelcolor="#525050")
    # Reducir cantidad de ticks X
    n = len(x)
    if n > 12:
        step = max(1, n // 8)
        for i, lbl_t in enumerate(ax.get_xticklabels()):
            if i % step != 0:
                lbl_t.set_visible(False)
    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", facecolor="white", dpi=140)
    plt.close(fig)
    buf.seek(0)
    return buf

def bar_chart(arr, x_key, y_key, title=None, y_fmt=None, n_meses=24, semantic=True):
    """Chart de barras con colores semánticos (verde positivo / rojo negativo)."""
    if not arr:
        return None
    data = arr[-n_meses:]
    x = [r.get(x_key, "") for r in data]
    y = [r.get(y_key) or 0 for r in data]
    colors = [VERDE_HX if v >= 0 else ROJO_HX for v in y] if semantic else [CYAN_HX] * len(y)
    fig, ax = plt.subplots(figsize=(6.5, 3.4), dpi=140)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.bar(x, y, color=colors, width=0.7, edgecolor="none")
    ax.axhline(0, color="#9A9990", linewidth=0.5)
    style_axes(ax, title)
    if y_fmt:
        ax.yaxis.set_major_formatter(FuncFormatter(y_fmt))
    n = len(x)
    if n > 12:
        step = max(1, n // 8)
        for i, lbl_t in enumerate(ax.get_xticklabels()):
            if i % step != 0:
                lbl_t.set_visible(False)
    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight", facecolor="white", dpi=140)
    plt.close(fig)
    buf.seek(0)
    return buf

# ──────────────────────────────────────────────────────────────────────────────
# Templates por serie — generan: título, lede, bullets, chart, tabla
# ──────────────────────────────────────────────────────────────────────────────
def tpl_ipc(p, D):
    d = p["datos"]
    arr = D.get("ipc", [])
    arr_real = [r for r in arr if not r.get("proj") and r.get("vm") is not None]
    titulo = f"INFLACIÓN {lbl_mes(p['periodo']).upper()}: {fmt_pct(d.get('vm'))}"
    lede = (f"El IPC nivel general aumentó {fmt_pct(d.get('vm'))} mensual en {lbl_mes(p['periodo'])}, "
            f"acumulando {d.get('via','—')}% interanual. INDEC publicó el indicador "
            f"con la metodología vigente.")
    bullets = [
        f"Variación mensual: {fmt_pct(d.get('vm'))}",
        f"Variación interanual: {d.get('via','—')}%",
        f"Acumulado últimos 6 meses: {sum((r.get('vm') or 0) for r in arr_real[-6:]):.1f}% (suma simple)",
    ]
    chart = bar_chart(arr_real, "f", "vm",
                     title="Inflación mensual %",
                     y_fmt=lambda x, p: f"{x:.0f}%",
                     n_meses=24)
    # Tabla últimos 6 meses
    tabla = [["Mes", "VM %", "VIA %"]]
    for r in arr_real[-6:][::-1]:
        tabla.append([lbl_mes(r["f"]).capitalize(), fmt_pct(r.get("vm")), fmt_pct(r.get("via"))])
    return titulo, lede, bullets, chart, tabla

def tpl_ipim(p, D):
    d = p["datos"]
    arr = D.get(p["serie"], [])
    arr_v = [r for r in arr if r.get("vm") is not None]
    nombre = {"ipim":"IPIM mayorista", "ipib":"IPIB básicos", "ipp":"IPP productor"}.get(p["serie"], "IPIM")
    titulo = f"{nombre.upper()} {lbl_mes(p['periodo']).upper()}: +{d.get('vm','—')}%"
    lede = (f"El {nombre} registró un aumento de +{d.get('vm','—')}% mensual en "
            f"{lbl_mes(p['periodo'])}, acumulando {d.get('via','—')}% interanual. "
            f"Suele anticipar el IPC minorista.")
    bullets = [
        f"Nivel general: +{d.get('vm','—')}% mensual",
        f"Interanual: {d.get('via','—')}%",
    ]
    # Comparación con IPC del mismo mes
    ipc = next((r for r in D.get("ipc", []) if r.get("f") == p["periodo"] and not r.get("proj")), None)
    if ipc and ipc.get("vm") is not None:
        br = d.get("vm", 0) - ipc["vm"]
        bullets.append(f"Brecha vs IPC minorista: {fmt_pct(br)} pp (IPC {fmt_pct(ipc['vm'])})")
    chart = bar_chart(arr_v, "f", "vm",
                     title=f"{nombre} mensual %",
                     y_fmt=lambda x, p: f"{x:.0f}%", n_meses=24)
    tabla = [["Mes", "VM %", "VIA %"]]
    for r in arr_v[-6:][::-1]:
        tabla.append([lbl_mes(r["f"]).capitalize(), fmt_pct(r.get("vm"), 2), fmt_pct(r.get("via"))])
    return titulo, lede, bullets, chart, tabla

def tpl_bc(p, D):
    d = p["datos"]
    arr = D.get("bc", [])
    arr_v = [r for r in arr if r.get("saldo") is not None]
    saldo = d.get("saldo", 0)
    expo = d.get("expo", 0)
    impo = d.get("impo_abs") or abs(d.get("impo", 0))
    signo = "SUPERÁVIT" if saldo >= 0 else "DÉFICIT"
    titulo = f"BALANZA COMERCIAL {lbl_mes(p['periodo']).upper()}: {signo} USD {fmt_n(abs(saldo))}M"
    # Var i.a. de expo/impo
    prev_y_period = f"{int(p['periodo'][:4]) - 1}{p['periodo'][4:]}"
    prev_y = next((r for r in arr if r.get("f") == prev_y_period), None)
    via_e = f"{((expo / prev_y['expo'] - 1) * 100):+.1f}%" if prev_y and prev_y.get("expo") else "—"
    via_i = f"{((impo / prev_y['impo_abs'] - 1) * 100):+.1f}%" if prev_y and prev_y.get("impo_abs") else "—"
    lede = (f"El intercambio comercial argentino en {lbl_mes(p['periodo'])} arrojó un "
            f"{signo.lower()} de USD {fmt_n(abs(saldo))} millones. "
            f"Exportaciones: USD {fmt_n(expo)} M ({via_e} i.a.). "
            f"Importaciones: USD {fmt_n(impo)} M ({via_i} i.a.).")
    bullets = [
        f"Exportaciones: USD {fmt_n(expo)} millones ({via_e} i.a.)",
        f"Importaciones: USD {fmt_n(impo)} millones ({via_i} i.a.)",
        f"Saldo comercial: USD {('+' if saldo>=0 else '')}{fmt_n(saldo)} M",
    ]
    chart = bar_chart(arr_v, "f", "saldo",
                     title="Saldo comercial mensual (USD millones)",
                     y_fmt=lambda x, p: f"{x:,.0f}".replace(",", "."), n_meses=24)
    tabla = [["Mes", "Expo USD M", "Impo USD M", "Saldo USD M"]]
    for r in arr_v[-6:][::-1]:
        tabla.append([lbl_mes(r["f"]).capitalize(), fmt_n(r.get("expo")),
                     fmt_n(r.get("impo_abs") or abs(r.get("impo", 0))), fmt_n(r.get("saldo"))])
    return titulo, lede, bullets, chart, tabla

def tpl_fiscal(p, D):
    d = p["datos"]
    arr = D.get("fiscal", [])
    arr_v = [r for r in arr if r.get("sal_prim") is not None]
    sp = d.get("sal_prim", 0)
    sf = d.get("sal_fin", 0)
    signo = "SUPERÁVIT" if sp >= 0 else "DÉFICIT"
    titulo = f"FISCAL {lbl_mes(p['periodo']).upper()}: {signo} PRIMARIO ${fmt_n(abs(sp/1000))}K M$"
    lede = (f"El Sector Público Nacional No Financiero registró un resultado primario de "
            f"{('+' if sp>=0 else '')}{fmt_n(sp/1000)} mil millones de pesos en {lbl_mes(p['periodo'])}. "
            f"El resultado financiero alcanzó {('+' if sf>=0 else '')}{fmt_n(sf/1000)}K M$.")
    bullets = [
        f"Resultado primario: {('+' if sp>=0 else '')}{fmt_n(sp/1000)} mil millones de pesos",
        f"Resultado financiero: {('+' if sf>=0 else '')}{fmt_n(sf/1000)} mil millones",
        f"Fuente: Secretaría de Hacienda, base caja (Metodología 2017)",
    ]
    chart = bar_chart(arr_v, "f", "sal_prim",
                     title="Resultado primario SPNF (millones de pesos)",
                     y_fmt=lambda x, p: f"{x/1000:,.0f}K".replace(",", "."), n_meses=24)
    tabla = [["Mes", "Primario M$", "Financiero M$"]]
    for r in arr_v[-6:][::-1]:
        tabla.append([lbl_mes(r["f"]).capitalize(),
                     fmt_n(r.get("sal_prim")), fmt_n(r.get("sal_fin"))])
    return titulo, lede, bullets, chart, tabla

def tpl_emae(p, D):
    d = p["datos"]
    arr = D.get("emae", [])
    arr_v = [r for r in arr if r.get("via") is not None]
    via = d.get("via", 0)
    vm = d.get("vm")
    verbo = "CRECIÓ" if via >= 0 else "CAYÓ"
    titulo = f"EMAE {lbl_mes(p['periodo']).upper()}: ACTIVIDAD {verbo} {fmt_pct(via)} I.A."
    lede = (f"El Estimador Mensual de Actividad Económica (EMAE) registró una variación interanual "
            f"de {fmt_pct(via)} en {lbl_mes(p['periodo'])}."
            + (f" En la medición desestacionalizada, el indicador se movió {fmt_pct(vm)} respecto al mes anterior." if vm is not None else ""))
    bullets = [
        f"Variación interanual: {fmt_pct(via)}",
    ]
    if vm is not None:
        bullets.append(f"Variación mensual desestacionalizada: {fmt_pct(vm)}")
    bullets.append("Serie INDEC, base 2004=100. Mide el nivel general de actividad.")
    chart = bar_chart(arr_v, "f", "via",
                     title="EMAE variación interanual %",
                     y_fmt=lambda x, p: f"{x:.0f}%", n_meses=24)
    tabla = [["Mes", "Original", "VIA %", "Desest."]]
    for r in arr_v[-6:][::-1]:
        tabla.append([lbl_mes(r["f"]).capitalize(),
                     fmt_n(r.get("orig"), 1), fmt_pct(r.get("via")),
                     fmt_n(r.get("dest"), 1) if r.get("dest") else "—"])
    return titulo, lede, bullets, chart, tabla

def tpl_ipi(p, D):
    d = p["datos"]
    arr = D.get("ipi", [])
    arr_v = [r for r in arr if r.get("via") is not None]
    via = d.get("via", 0)
    verbo = "CRECIÓ" if via >= 0 else "CAYÓ"
    titulo = f"INDUSTRIA {lbl_mes(p['periodo']).upper()}: IPI {verbo} {fmt_pct(via)} I.A."
    lede = (f"El Índice de Producción Industrial Manufacturero (IPI) registró una variación "
            f"interanual de {fmt_pct(via)} en {lbl_mes(p['periodo'])}.")
    bullets = [
        f"Variación interanual: {fmt_pct(via)}",
        "Serie INDEC, base 2004=100. Mide la producción manufacturera (16 ramas).",
    ]
    chart = bar_chart(arr_v, "f", "via",
                     title="IPI variación interanual %",
                     y_fmt=lambda x, p: f"{x:.0f}%", n_meses=24)
    tabla = [["Mes", "Original", "VIA %"]]
    for r in arr_v[-6:][::-1]:
        tabla.append([lbl_mes(r["f"]).capitalize(),
                     fmt_n(r.get("orig"), 1), fmt_pct(r.get("via"))])
    return titulo, lede, bullets, chart, tabla

def tpl_super(p, D):
    d = p["datos"]
    arr = D.get("super", [])
    arr_v = [r for r in arr if r.get("via_real") is not None]
    via = d.get("via_real", 0)
    vm = d.get("vm_dest")
    verbo = "CRECIERON" if via >= 0 else "CAYERON"
    titulo = f"SUPERMERCADOS {lbl_mes(p['periodo']).upper()}: VENTAS REALES {verbo} {fmt_pct(via)}"
    lede = (f"Las ventas de supermercados a precios constantes registraron una variación interanual real "
            f"de {fmt_pct(via)} en {lbl_mes(p['periodo'])}."
            + (f" Desestacionalizada: {fmt_pct(vm)} m/m." if vm is not None else ""))
    bullets = [
        f"Variación interanual real: {fmt_pct(via)}",
    ]
    if vm is not None:
        bullets.append(f"Variación mensual desestacionalizada: {fmt_pct(vm)}")
    if d.get("acum_real") is not None:
        bullets.append(f"Acumulado del año: {fmt_pct(d['acum_real'])}")
    chart = bar_chart(arr_v, "f", "via_real",
                     title="Ventas supermercados — variación i.a. real %",
                     y_fmt=lambda x, p: f"{x:.0f}%", n_meses=24)
    tabla = [["Mes", "Índice real", "VIA real %", "Acumulado %"]]
    for r in arr_v[-6:][::-1]:
        tabla.append([lbl_mes(r["f"]).capitalize(),
                     fmt_n(r.get("n_real"), 1), fmt_pct(r.get("via_real")),
                     fmt_pct(r.get("acum_real"))])
    return titulo, lede, bullets, chart, tabla

def tpl_salarios(p, D):
    d = p["datos"]
    arr = D.get("salarios", [])
    arr_v = [r for r in arr if r.get("is_r") is not None]
    titulo = f"SALARIOS {lbl_mes(p['periodo']).upper()}: ÍNDICE TOTAL REGISTRADO {fmt_n(d.get('is_r',0))}"
    lede = (f"El Índice de Salarios INDEC alcanzó {fmt_n(d.get('is_r',0))} para el sector registrado en "
            f"{lbl_mes(p['periodo'])}. Privado registrado: {fmt_n(d.get('real_priv',0))}. "
            f"Público: {fmt_n(d.get('real_pub',0))}.")
    bullets = [
        f"Total registrado: {fmt_n(d.get('is_r','—'))}",
        f"Privado registrado: {fmt_n(d.get('real_priv','—'))}",
        f"Sector público: {fmt_n(d.get('real_pub','—'))}",
    ]
    if d.get("no_reg") is not None:
        bullets.append(f"No registrado: {fmt_n(d['no_reg'])}")
    chart = line_chart(arr_v, "f", ["is_r", "real_priv", "real_pub"],
                       ["Total registrado", "Privado", "Público"],
                       [CYAN_HX, "#1E4A8A", "#A85C1E"],
                       title="Índice de salarios (base oct 2016=100)",
                       n_meses=24)
    tabla = [["Mes", "Total reg.", "Privado", "Público"]]
    for r in arr_v[-6:][::-1]:
        tabla.append([lbl_mes(r["f"]).capitalize(),
                     fmt_n(r.get("is_r")), fmt_n(r.get("real_priv")), fmt_n(r.get("real_pub"))])
    return titulo, lede, bullets, chart, tabla

def tpl_mora(p, D):
    d = p["datos"]
    arr = D.get("mora", [])
    arr_v = [r for r in arr if r.get("fam") is not None]
    rec_per = next((r for r in arr if r.get("f") == p["periodo"]), {})
    is_prov = bool(rec_per.get("prov"))
    fuente_str = rec_per.get("fuente", "BCRA Informe sobre Bancos")
    fam = d.get("fam", 0)
    emp = d.get("emp", 0)
    prov_tag = " (PRELIMINAR)" if is_prov else ""
    titulo = f"MORA BANCARIA {lbl_mes(p['periodo']).upper()}: FAMILIAS {fam}% · EMPRESAS {emp}%{prov_tag}"
    # Delta vs mes anterior
    prev = next((r for r in reversed(arr_v) if r.get("f") < p["periodo"]), None)
    delta_str = ""
    if prev and prev.get("fam") is not None:
        delta = fam - prev["fam"]
        signo = "+" if delta >= 0 else ""
        delta_str = f" (venía de {prev['fam']:.2f}% en {lbl_mes(prev['f'])}, {signo}{delta:.2f} pp)"
    lede = (f"La irregularidad de cartera alcanzó {fam}% en familias y {emp}% en empresas "
            f"en {lbl_mes(p['periodo'])}{delta_str}.")
    if is_prov:
        lede += f" Dato preliminar reportado por {fuente_str}; el BCRA publicará el oficial en su próximo Informe sobre Bancos."
    bullets = [
        f"Familias: {fam}% del total de préstamos en situación irregular",
        f"Empresas: {emp}%",
        f"Fuente: {fuente_str}",
    ]
    if is_prov:
        bullets.append("⚠ Dato preliminar — sujeto a revisión cuando el BCRA publique el Informe oficial")
    chart = line_chart(arr_v, "f", ["fam", "emp"],
                       ["Familias", "Empresas"],
                       [ROJO_HX, "#A85C1E"],
                       title="Irregularidad de cartera %",
                       y_fmt=lambda x, p: f"{x:.0f}%", n_meses=24)
    tabla = [["Mes", "Familias %", "Empresas %"]]
    for r in arr_v[-6:][::-1]:
        flag = " *" if r.get("prov") else ""
        tabla.append([lbl_mes(r["f"]).capitalize() + flag,
                     f"{r.get('fam')}%", f"{r.get('emp')}%"])
    return titulo, lede, bullets, chart, tabla

def tpl_rec(p, D):
    """Recaudación tributaria con cálculo de variación real i.a. (usa IPC real o proyección REM)."""
    d = p["datos"]
    per = p["periodo"]
    arr = D.get("rec", [])
    arr_v = [r for r in arr if r.get("tot") is not None]
    tot = d.get("tot", 0)
    # i.a. nominal y real
    prev_y = next((r for r in arr if r.get("f") == f"{int(per[:4])-1}-{per[5:7]}"), None)
    via_nom = None
    via_real = None
    via_real_src = ""
    if prev_y and prev_y.get("tot"):
        via_nom = (tot / prev_y["tot"] - 1) * 100
        ipc_act = next((r for r in D.get("ipc", []) if r.get("f") == per), None)
        ipc_prev = next((r for r in D.get("ipc", []) if r.get("f") == f"{int(per[:4])-1}-{per[5:7]}"), None)
        if ipc_act and ipc_act.get("n") and ipc_prev and ipc_prev.get("n"):
            infl_ia = (ipc_act["n"] / ipc_prev["n"] - 1) * 100
            via_real = ((1 + via_nom / 100) / (1 + infl_ia / 100) - 1) * 100
            via_real_src = "REM" if ipc_act.get("proj") else "IPC oficial"
    # Headline
    if via_real is not None:
        verbo = "creció" if via_real >= 0 else "cayó"
        titulo = f"RECAUDACIÓN {lbl_mes(per).upper()}: {verbo.upper()} {via_real:+.1f}% REAL I.A."
    elif via_nom is not None:
        titulo = f"RECAUDACIÓN {lbl_mes(per).upper()}: {via_nom:+.1f}% NOMINAL I.A."
    else:
        titulo = f"RECAUDACIÓN {lbl_mes(per).upper()}: ${fmt_n(tot/1000)}K M$"
    lede_parts = [f"AFIP/ARCA recaudó ${fmt_n(tot/1000)}K millones en {lbl_mes(per)}."]
    if via_nom is not None:
        lede_parts.append(f"Variación nominal interanual: {via_nom:+.1f}%.")
    if via_real is not None:
        lede_parts.append(f"En términos reales (deflactado por {via_real_src}): {via_real:+.1f}% i.a.")
    # Detectar quiebre de tendencia
    serie_real = [r for r in arr_v if r.get("via_r") is not None and r["f"] < per][-6:]
    negs = sum(1 for r in serie_real if r.get("via_r", 0) < 0)
    if via_real is not None and via_real >= 0 and negs >= 3:
        lede_parts.append(f"Quiebre de tendencia: venía con {negs} meses de caída real interanual.")
    lede = " ".join(lede_parts)
    bullets = [f"Total: ${fmt_n(tot/1000)}K M$"]
    if via_nom is not None:
        bullets.append(f"Variación nominal i.a.: {via_nom:+.1f}%")
    if via_real is not None:
        bullets.append(f"Variación real i.a.: {via_real:+.1f}% (deflactor: {via_real_src})")
    if d.get("iva") is not None:
        bullets.append(f"IVA: ${fmt_n(d['iva']/1000)}K M$")
    if d.get("gan") is not None:
        bullets.append(f"Ganancias: ${fmt_n(d['gan']/1000)}K M$")
    bullets.append("Fuente: ARCA (ex AFIP) · indec.gob.ar series")
    # Chart: var real i.a. últimos 24m
    arr_chart = [r for r in arr_v if r.get("via_r") is not None][-24:]
    chart = bar_chart(arr_chart, "f", "via_r",
                      title="Recaudación tributaria · variación real interanual %",
                      y_fmt=lambda x, p: f"{x:+.0f}%",
                      n_meses=24, semantic=True) if arr_chart else None
    # Tabla
    tabla = [["Mes", "Total M$", "Var. nom. i.a.", "Var. real i.a."]]
    for r in arr_v[-6:][::-1]:
        tabla.append([
            lbl_mes(r["f"]).capitalize(),
            fmt_n(r.get("tot", 0)/1000) + "K",
            f"{r.get('via','—')}%" if r.get("via") is not None else "—",
            f"{r.get('via_r','—')}%" if r.get("via_r") is not None else "—",
        ])
    return titulo, lede, bullets, chart, tabla

def tpl_turismo(p, D):
    d = p["datos"]
    per = p["periodo"]
    arr = D.get("turismo", [])
    arr_v = [r for r in arr if r.get("sal") is not None]
    rec = d.get("rec", 0)
    em = d.get("em", 0)
    sal = d.get("sal", rec - em)
    titulo = f"TURISMO INTERNACIONAL {lbl_mes(per).upper()}: SALDO {fmt_n(sal)} PERSONAS"
    # Var i.a.
    prev_y = next((r for r in arr if r.get("f") == f"{int(per[:4])-1}-{per[5:7]}"), None)
    via_str = ""
    if prev_y and prev_y.get("rec") and prev_y.get("em"):
        via_rec = (rec / prev_y["rec"] - 1) * 100
        via_em = (em / prev_y["em"] - 1) * 100
        via_str = f" Receptivo {via_rec:+.1f}% i.a., emisivo {via_em:+.1f}% i.a."
    lede = (f"En {lbl_mes(per)} ingresaron {fmt_n(rec)} turistas no residentes y "
            f"salieron {fmt_n(em)} argentinos al exterior, con saldo de {fmt_n(sal)} personas.{via_str}")
    bullets = [
        f"Receptivo: {fmt_n(rec)} personas",
        f"Emisivo: {fmt_n(em)} personas",
        f"Saldo (receptivo − emisivo): {fmt_n(sal)} personas",
        "Fuente: INDEC, Estadísticas de Turismo Internacional (todas las vías)",
    ]
    chart = line_chart(arr_v, "f", ["rec", "em"],
                       ["Receptivo", "Emisivo"],
                       [CYAN_HX, ROJO_HX],
                       title="Turismo internacional · personas por mes",
                       n_meses=24)
    tabla = [["Mes", "Receptivo", "Emisivo", "Saldo"]]
    for r in arr_v[-6:][::-1]:
        tabla.append([lbl_mes(r["f"]).capitalize(),
                     fmt_n(r.get("rec", 0)), fmt_n(r.get("em", 0)), fmt_n(r.get("sal", 0))])
    return titulo, lede, bullets, chart, tabla

def tpl_default(p, D):
    d = p["datos"]
    titulo = f"{p['label'].upper()} {lbl_mes(p['periodo']).upper()}"
    lede = f"Se publicó la actualización de {p['label']} para {lbl_mes(p['periodo'])}. Fuente: {p.get('fuente','INDEC')}."
    bullets = [f"{k}: {fmt_n(v, 2)}" for k, v in list(d.items())[:5]]
    return titulo, lede, bullets, None, None

TEMPLATES = {
    "ipc": tpl_ipc,
    "ipim": tpl_ipim, "ipib": tpl_ipim, "ipp": tpl_ipim,
    "bc": tpl_bc,
    "fiscal": tpl_fiscal,
    "emae": tpl_emae,
    "ipi": tpl_ipi, "isac": tpl_ipi, "uci": tpl_ipi,
    "super": tpl_super, "mayor": tpl_super,
    "salarios": tpl_salarios, "ripte": tpl_salarios,
    "mora": tpl_mora,
    "rec": tpl_rec,
    "turismo": tpl_turismo,
}

# ──────────────────────────────────────────────────────────────────────────────
# Construcción del slide
# ──────────────────────────────────────────────────────────────────────────────
def add_text_run(paragraph, text, *, bold=False, italic=False, size=14, color=None, font=FONT):
    r = paragraph.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    if color is not None:
        r.font.color.rgb = color
    return r

def style_table(table, header_rows=1):
    """Aplica estilo Comex a una tabla PPTX."""
    n_cols = len(table.columns)
    n_rows = len(table.rows)
    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            tf = cell.text_frame
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(10)
                    if r < header_rows:
                        run.font.bold = True
                        run.font.color.rgb = BLANCO
                    else:
                        run.font.color.rgb = GRIS_TX
            # Background
            from pptx.oxml.ns import qn as _qn
            tcPr = cell._tc.get_or_add_tcPr()
            for fill_old in tcPr.findall(_qn("a:solidFill")):
                tcPr.remove(fill_old)
            solid = etree.SubElement(tcPr, _qn("a:solidFill"))
            clr = etree.SubElement(solid, _qn("a:srgbClr"))
            if r < header_rows:
                clr.set("val", "00B2C9")
            else:
                clr.set("val", "FFFFFF" if r % 2 == 1 else "F8F8F6")

def build_slide(prs, pub, D):
    serie = pub.get("serie", "")
    tpl = TEMPLATES.get(serie, tpl_default)
    titulo, lede, bullets, chart_buf, tabla = tpl(pub, D)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Banner cyan
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.3))
    banner.fill.solid()
    banner.fill.fore_color.rgb = CYAN
    banner.line.fill.background()

    # Título
    tit_l, tit_t, tit_w = Inches(0.45), Inches(0.5), prs.slide_width - Inches(0.9)
    tit_box = slide.shapes.add_textbox(tit_l, tit_t, tit_w, Inches(1.0))
    tf = tit_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    add_text_run(tf.paragraphs[0], titulo, bold=True, size=26, color=CYAN)
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Sub: fuente + periodo
    sub_box = slide.shapes.add_textbox(tit_l, Inches(1.4), tit_w, Inches(0.3))
    add_text_run(sub_box.text_frame.paragraphs[0],
                 f"{pub.get('fuente','INDEC')} · período {pub.get('periodo','')}",
                 italic=True, size=11, color=GRIS_LX)

    # Layout:
    # Columna izquierda (lede + bullets): x=0.45, w=6.0
    # Columna derecha (chart + tabla):    x=6.65, w=6.2
    col_l_w = Inches(6.0)
    col_r_x = Inches(6.65)
    col_r_w = Inches(6.2)

    # Lede izquierda
    lede_box = slide.shapes.add_textbox(tit_l, Inches(1.85), col_l_w, Inches(1.8))
    ltf = lede_box.text_frame
    ltf.word_wrap = True
    ltf.margin_left = ltf.margin_right = Emu(0)
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.JUSTIFY
    add_text_run(lp, lede, size=13, color=GRIS_TX)

    # Bullets izquierda
    bul_box = slide.shapes.add_textbox(tit_l, Inches(3.85), col_l_w, Inches(3.0))
    btf = bul_box.text_frame
    btf.word_wrap = True
    btf.margin_left = btf.margin_right = Emu(0)
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        # Bullet en cyan, texto en gris
        add_text_run(p, "●  ", size=13, color=CYAN, bold=True)
        add_text_run(p, b, size=13, color=GRIS_TX)

    # Chart derecha (matplotlib PNG)
    if chart_buf is not None:
        slide.shapes.add_picture(chart_buf, col_r_x, Inches(1.85),
                                 width=col_r_w, height=Inches(2.8))

    # Tabla derecha abajo
    if tabla and len(tabla) > 1:
        n_rows = len(tabla)
        n_cols = len(tabla[0])
        tbl_top = Inches(4.85)
        tbl_h = Inches(0.32 * n_rows)
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, col_r_x, tbl_top, col_r_w, tbl_h)
        tbl = tbl_shape.table
        for r_idx, row in enumerate(tabla):
            for c_idx, val in enumerate(row):
                cell = tbl.cell(r_idx, c_idx)
                cell.text = str(val)
        style_table(tbl, header_rows=1)

    # Footer
    foot_box = slide.shapes.add_textbox(tit_l, prs.slide_height - Inches(0.4),
                                        prs.slide_width - Inches(0.9), Inches(0.3))
    add_text_run(foot_box.text_frame.paragraphs[0],
                 f"Fuente: elaboración propia en base a {pub.get('fuente','INDEC')}",
                 italic=True, size=10, color=GRIS_LX)

def build_cover(prs, n_pubs, lunes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = CYAN
    fondo.line.fill.background()

    # Decoración: línea blanca a la izquierda
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.6), Inches(0.06), Inches(3.0))
    deco.fill.solid()
    deco.fill.fore_color.rgb = BLANCO
    deco.line.fill.background()

    tit = slide.shapes.add_textbox(Inches(0.95), Inches(2.5), prs.slide_width - Inches(1.5), Inches(1.4))
    p = tit.text_frame.paragraphs[0]
    add_text_run(p, "PANORAMA MACRO", bold=True, size=54, color=BLANCO)

    sub = slide.shapes.add_textbox(Inches(0.95), Inches(3.8), prs.slide_width - Inches(1.5), Inches(0.8))
    p2 = sub.text_frame.paragraphs[0]
    add_text_run(p2, f"Publicaciones de la semana del {lunes.strftime('%d de ')}{MESES_ES[lunes.month - 1]}{lunes.strftime(' de %Y')}",
                 size=22, color=BLANCO)

    sub2 = slide.shapes.add_textbox(Inches(0.95), Inches(4.6), prs.slide_width - Inches(1.5), Inches(0.5))
    add_text_run(sub2.text_frame.paragraphs[0],
                 f"{n_pubs} publicaciones · datos oficiales INDEC / BCRA / Hacienda",
                 size=14, italic=True, color=BLANCO)

def main():
    data = json.loads(DATA_PATH.read_text(encoding="utf-8"))
    pubs = (data.get("_meta") or {}).get("publicaciones") or []
    if not pubs:
        print("Sin publicaciones tracked. No genero PPTX.")
        return 1
    lunes = ultimo_lunes()
    lunes_iso = lunes.isoformat()
    semana = [p for p in pubs if (p.get("publicado_at") or p.get("detectado_at","")) >= lunes_iso]
    if not semana:
        print("Sin publicaciones esta semana. No genero PPTX.")
        return 1
    semana.sort(key=lambda p: p.get("publicado_at") or p.get("detectado_at",""))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_cover(prs, len(semana), lunes)
    for pub in semana:
        try:
            build_slide(prs, pub, data)
        except Exception as e:
            print(f"Error en slide {pub.get('serie')} {pub.get('periodo')}: {e}")
            import traceback
            traceback.print_exc()

    prs.save(str(OUT_PATH))
    print(f"OK — {OUT_PATH} ({OUT_PATH.stat().st_size/1024:.1f} KB, {len(semana)+1} slides)")
    return 0

if __name__ == "__main__":
    raise SystemExit(main())
